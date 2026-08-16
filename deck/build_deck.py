# -*- coding: utf-8 -*-
"""SKN 피치덱 v3 — Pretendard · 실제 로고/오브 · 임베드 차트 · 방탄 BM. 15장 16:9."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import os

BASE = "/private/tmp/claude-501/-Users-xhae000-orca-projects-central-hack/bb8256fd-8a12-465c-b443-4aa9d54c15c8/scratchpad"
CH = f"{BASE}/charts"; AS = f"{BASE}/assets"

# 팔레트
PAPER=RGBColor(0xFA,0xF8,0xF3); PANEL=RGBColor(0xF1,0xED,0xE5); INK=RGBColor(0x1A,0x17,0x12)
SUB=RGBColor(0x6B,0x64,0x59); FAINT=RGBColor(0xA7,0x9E,0x90); HAIR=RGBColor(0xDD,0xD6,0xC8)
ACCENT=RGBColor(0x4E,0x5B,0x43); CLAY=RGBColor(0x9A,0x6A,0x4B)
INKBG=RGBColor(0x18,0x15,0x0F); PAPERON=RGBColor(0xEE,0xE9,0xDE); SUBON=RGBColor(0x9A,0x93,0x84)
HAIRON=RGBColor(0x3A,0x35,0x2B); SAGEON=RGBColor(0x9C,0xAB,0x8C)

# 폰트(웨이트=패밀리)
TH="Pretendard Thin"; XL="Pretendard ExtraLight"; LT="Pretendard Light"; RG="Pretendard"
MD="Pretendard Medium"; SB="Pretendard SemiBold"; XB="Pretendard ExtraBold"; BK="Pretendard Black"

PW,PH=13.333,7.5; ML,MR=0.95,0.95; CW=PW-ML-MR
prs=Presentation(); prs.slide_width=Inches(PW); prs.slide_height=Inches(PH)
BLANK=prs.slide_layouts[6]

def slide(bg=PAPER):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    sp=r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp); return s

def _spc(run,p): run._r.get_or_add_rPr().set('spc',str(int(p*100)))

def T(s,x,y,w,h,runs,al=PP_ALIGN.LEFT,an=MSO_ANCHOR.TOP,ls=1.0,sa=0.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=an
    for m in('left','right','top','bottom'): setattr(tf,'margin_'+m,0)
    bp=tf._txBody.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit','a:spAutoFit'):
        e=bp.find(qn(tag));
        if e is not None: bp.remove(e)
    bp.append(tf._txBody.makeelement(qn('a:noAutofit'),{}))
    paras=runs if isinstance(runs[0],list) else [runs]
    for pi,pa in enumerate(paras):
        p=tf.paragraphs[0] if pi==0 else tf.add_paragraph()
        p.alignment=al; p.line_spacing=ls; p.space_before=Pt(0)
        if sa: p.space_after=Pt(sa)
        for (t,sz,c,f,tr) in pa:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.color.rgb=c; r.font.name=f
            if tr: _spc(r,tr)
    return tb

def R(t,sz,c=INK,f=RG,tr=0.0): return (t,sz,c,f,tr)

def hline(s,x,y,w,c=HAIR,wt=1.0):
    ln=s.shapes.add_connector(2,Inches(x),Inches(y),Inches(x+w),Inches(y))
    ln.line.color.rgb=c; ln.line.width=Pt(wt); ln.shadow.inherit=False; return ln
def vline(s,x,y,h,c=HAIR,wt=1.0):
    ln=s.shapes.add_connector(2,Inches(x),Inches(y),Inches(x),Inches(y+h))
    ln.line.color.rgb=c; ln.line.width=Pt(wt); ln.shadow.inherit=False; return ln
def orb_ring(s,x,y,d,c=HAIR,wt=1.1):
    o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    o.fill.background(); o.line.color.rgb=c; o.line.width=Pt(wt); o.shadow.inherit=False; return o
def dot(s,x,y,d,c=ACCENT):
    o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb=c; o.line.fill.background(); o.shadow.inherit=False; return o

_ratio={}
def ratio(path):
    if path not in _ratio:
        w,h=Image.open(path).size; _ratio[path]=w/h
    return _ratio[path]
def pic_w(s,path,x,y,w):
    p=s.shapes.add_picture(path,Inches(x),Inches(y),width=Inches(w)); p.shadow.inherit=False; return p
def pic_h(s,path,x,y,h):
    p=s.shapes.add_picture(path,Inches(x),Inches(y),height=Inches(h)); p.shadow.inherit=False; return p

def eyebrow(s,num,label,y=0.82,dark=False):
    T(s,ML,y,8,0.34,[R(num,12.5,(SAGEON if dark else ACCENT),SB),
                     R("   "+label,10.5,(SUBON if dark else FAINT),MD,3.0)],an=MSO_ANCHOR.MIDDLE)
def foot(s,idx,dark=False):
    pic_w(s,f"{AS}/logo_paper.png" if dark else f"{AS}/logo_ink.png",ML,PH-0.62,0.62)
    T(s,PW-MR-1.0,PH-0.6,1.0,0.3,[R(f"{idx:02d}",11.5,(SUBON if dark else FAINT),RG)],
      al=PP_ALIGN.RIGHT,an=MSO_ANCHOR.MIDDLE)
def source(s,txt):
    T(s,ML+0.85,PH-0.6,CW-1.85,0.3,[R("출처  ",9,FAINT,MD),R(txt,9,FAINT,RG)],an=MSO_ANCHOR.MIDDLE)

def phone(s,cx,top,h,label,sub):
    w=h*(300.0/620.0); x=cx-w/2
    fr=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(top),Inches(w),Inches(h))
    fr.adjustments[0]=0.06; fr.fill.solid(); fr.fill.fore_color.rgb=PANEL
    fr.line.color.rgb=HAIR; fr.line.width=Pt(1.3); fr.shadow.inherit=False
    hline(s,x+w*0.37,top+0.26,w*0.26,c=HAIR,wt=2.0)
    T(s,x+0.2,top,w-0.4,h-0.3,[[R("화면",10,FAINT,MD,2.5)],[R(label,12.5,SUB,SB)],[R(sub,10,FAINT,RG)]],
      al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE,ls=1.3,sa=4)

# ============================================================
# 01 · TITLE
# ============================================================
s=slide(PAPER)
pic_w(s,f"{AS}/orb.png",8.35,1.35,4.5)
pic_w(s,f"{AS}/logo_ink.png",ML,0.8,1.5)
T(s,PW-MR-3.2,0.9,3.2,0.35,[R("AI · WELLNESS",11,FAINT,MD,3.0)],al=PP_ALIGN.RIGHT,an=MSO_ANCHOR.MIDDLE)
T(s,ML,2.5,9.2,3.0,
  [[R("써본 만큼,",50,INK,LT,-1.0)],
   [R("나를 더 잘 알게 되는",50,INK,LT,-1.0)],
   [R("스킨케어 경험 ",50,INK,LT,-1.0),R("아카이브",50,ACCENT,SB,-1.0)]],ls=1.18)
T(s,ML,5.6,8.6,0.8,
  [R("화장품을 지나가는 경험으로 두지 않습니다.  ",14.5,SUB,RG),
   R("AI가 내 기록에서 반복된 기준을 찾아 돌려줍니다.",14.5,INK,MD)],ls=1.45)
hline(s,ML,6.55,CW)
T(s,ML,6.72,8,0.35,[R("PERSONAL SKINCARE EXPERIENCE ARCHIVE",10.5,FAINT,MD,2.4)],an=MSO_ANCHOR.MIDDLE)
T(s,PW-MR-4,6.72,4,0.35,[R("skn.today",11,ACCENT,MD,0.3)],al=PP_ALIGN.RIGHT,an=MSO_ANCHOR.MIDDLE)

# ============================================================
# 02 · PROBLEM (chart)
# ============================================================
s=slide(PAPER)
eyebrow(s,"01","PROBLEM")
T(s,ML,1.4,11.5,1.4,
  [[R("코덕은 계속 탐색하지만,",35,INK,LT,-0.5)],
   [R("경험은 매번 ",35,INK,LT,-0.5),R("흩어진다",35,ACCENT,SB,-0.5)]],ls=1.18)
# 좌: 대형 숫자
T(s,ML,3.3,3.0,1.1,[R("9",60,INK,XL),R("개",21,SUB,RG)],an=MSO_ANCHOR.BOTTOM)
T(s,ML,4.45,3.2,0.4,[R("평균 보유 스킨케어 제품",12.5,SUB,MD)])
T(s,ML,4.95,3.0,1.1,[R("4",60,CLAY,XL),R("개",21,SUB,RG)],an=MSO_ANCHOR.BOTTOM)
T(s,ML,6.1,3.6,0.4,[R("다 쓰지 못하고 멈춘 제품",12.5,SUB,MD)])
vline(s,4.5,3.35,3.05,c=HAIR,wt=1.0)
# 우: 문제 바 차트
pic_w(s,f"{CH}/problem.png",5.0,3.7,6.9)
source(s,"SKN 예비 설문 n=214 · 20대 여성 · 2026.8 (수집 진행 중) · 값은 잠정")
foot(s,2)

# ============================================================
# 03 · THE GAP
# ============================================================
s=slide(PAPER)
eyebrow(s,"02","THE GAP")
T(s,ML,1.5,11.7,0.8,
  [R("정보는 많지만, ",33,INK,LT,-0.5),R("‘내 경험’",33,ACCENT,SB,-0.5),
   R("을 잇는 곳은 없다",33,INK,LT,-0.5)],ls=1.15)
rows=[("쇼핑몰 후기 · 커뮤니티","남이 남긴 제품 평가","내가 어떤 조건에서 무엇을 느꼈는지"),
      ("피부 일기 앱","날짜별 피부 상태","제품 · 루틴 · 사용감의 반복 패턴"),
      ("성분 분석 앱","정적인 성분 정보","내 실제 경험과 반대 기록"),
      ("범용 AI 채팅","지금 대화의 일반 설명","검증된 제품 사실 + 누적된 내 원본")]
ty=2.95; c1,c2=4.0,3.7
T(s,ML,ty,c1,0.3,[R("기존 대안",10.5,FAINT,MD,2.0)])
T(s,ML+c1,ty,c2,0.3,[R("주로 남는 것",10.5,FAINT,MD,2.0)])
T(s,ML+c1+c2,ty,CW-c1-c2,0.3,[R("SKN이 연결하는 것",10.5,ACCENT,MD,2.0)])
ry=ty+0.45; hline(s,ML,ry,CW,c=INK,wt=1.3); rh=0.82
for a,b,c in rows:
    T(s,ML,ry+0.02,c1-0.2,rh,[R(a,15,INK,SB)],an=MSO_ANCHOR.MIDDLE)
    T(s,ML+c1,ry+0.02,c2-0.2,rh,[R(b,13,SUB,RG)],an=MSO_ANCHOR.MIDDLE)
    T(s,ML+c1+c2,ry+0.02,CW-c1-c2,rh,[R(c,13.5,ACCENT,MD)],an=MSO_ANCHOR.MIDDLE)
    ry+=rh; hline(s,ML,ry,CW,c=HAIR,wt=1.0)
foot(s,3)

# ============================================================
# 04 · WHO (target donut)
# ============================================================
s=slide(PAPER)
eyebrow(s,"03","WHO")
T(s,ML,1.45,7.0,1.6,
  [[R("‘코덕’은 성별이 아니라",32,INK,LT,-0.5)],
   [R("행동",32,ACCENT,SB,-0.5),R("으로 정의됩니다",32,INK,LT,-0.5)]],ls=1.2)
beh=[("반복 탐색","새 제품을 반복해서 찾고 산다"),
     ("비교 사용","같은 카테고리를 여러 개 비교해 쓴다"),
     ("조건에 민감","제형·마무리·계절·시간대의 차이를 본다"),
     ("경험 재사용","과거 경험을 다음 선택에 쓰고 싶지만 흩어져 있다")]
ly=3.4; lw=6.9; hline(s,ML,ly,lw,c=INK,wt=1.3)
for i,(t,d) in enumerate(beh):
    T(s,ML,ly+0.03,0.8,0.7,[R(f"0{i+1}",16,ACCENT,LT)],an=MSO_ANCHOR.MIDDLE)
    T(s,ML+0.85,ly+0.03,2.1,0.7,[R(t,15,INK,SB)],an=MSO_ANCHOR.MIDDLE)
    T(s,ML+2.95,ly+0.03,lw-2.95,0.7,[R(d,12,SUB,RG)],an=MSO_ANCHOR.MIDDLE,ls=1.2)
    ly+=0.7; hline(s,ML,ly,lw,c=HAIR,wt=1.0)
# 우: 타깃 도넛
rx=ML+lw+0.6
pic_w(s,f"{CH}/target.png",9.0,3.25,2.85)
T(s,rx,6.35,PW-MR-rx,0.5,[R("20대 여성 중 ",12.5,SUB,RG),R("지속 탐색·구매형",12.5,INK,SB)],al=PP_ALIGN.CENTER)
source(s,"SKN 예비 설문 n=214 (수집 진행 중) · 값은 잠정")
foot(s,4)

# ============================================================
# 05 · VALIDATION (survey)
# ============================================================
s=slide(PAPER)
eyebrow(s,"04","VALIDATION")
T(s,ML,1.4,11.7,1.3,
  [R("예비 설문이 확인해준 것 — ",30,INK,LT,-0.5),
   R("문제도, 수요도 실재합니다",30,ACCENT,SB,-0.5)])
# 좌: 문제 실재 (recall 88%)
T(s,ML,2.75,3.4,0.35,[R("문제는 실재한다",12,FAINT,SB,1.5)])
pic_w(s,f"{CH}/recall.png",ML+0.15,3.1,2.9)
T(s,ML,6.15,3.3,0.6,[R("과거 경험을 기억·흩어진 메모로만 관리",11.5,SUB,RG)],ls=1.3)
vline(s,4.7,2.7,3.6,c=HAIR,wt=1.0)
# 우: 수요 실재 (intent + feature)
T(s,5.0,2.75,6,0.35,[R("수요도 실재한다",12,FAINT,SB,1.5)])
pic_w(s,f"{CH}/intent.png",5.0,3.2,4.05)
T(s,9.35,3.15,3.0,0.35,[R("가장 원하는 기능",11,FAINT,MD,1.0)])
pic_w(s,f"{CH}/feature.png",9.35,3.55,3.0)
source(s,"SKN 예비 설문 n=214 · 20대 여성 · 2026.8 (수집 진행 중) · 값은 잠정")
foot(s,5)

# ============================================================
# 06 · APPROACH (loop)
# ============================================================
s=slide(PAPER)
eyebrow(s,"05","APPROACH")
T(s,ML,1.4,11.8,1.4,
  [[R("쓴 경험을 제품·루틴에 연결하고,",31,INK,LT,-0.5)],
   [R("AI가 반복된 기준을 찾아 ",31,INK,LT,-0.5),R("다음 탐색에 되돌려줍니다",31,ACCENT,SB,-0.5)]],ls=1.2)
steps=[("01","DISCOVER","탐색","내 경험 기준으로\n후보를 본다"),
       ("02","EXPERIENCE","사용","실제 조합·순서로\n제품을 쓴다"),
       ("03","RECORD","기록","느낌을 내 말로\n남긴다"),
       ("04","CONNECT","연결","AI가 과거 기록과\n근거로 잇는다"),
       ("05","PATTERN","기준 발견","반복된 경험이\n내 기준이 된다")]
n=len(steps); ly=4.35; hline(s,ML,ly,CW,c=HAIR,wt=1.2); colw=CW/n
for i,(num,en,ko,d) in enumerate(steps):
    ac=(i==n-1); x0=ML+i*colw; cx=x0+colw/2
    dot(s,cx-0.05,ly-0.05,0.11,c=(ACCENT if ac else INK))
    T(s,x0,ly-1.1,colw-0.2,0.5,[R(num,15,ACCENT if ac else FAINT,LT)],al=PP_ALIGN.CENTER)
    T(s,x0,ly-0.62,colw-0.2,0.4,[R(ko,17,INK,SB)],al=PP_ALIGN.CENTER)
    T(s,x0,ly+0.26,colw-0.2,0.35,[R(en,9,ACCENT if ac else FAINT,MD,1.5)],al=PP_ALIGN.CENTER)
    T(s,x0,ly+0.66,colw-0.2,0.8,[R(d,11,SUB,RG)],al=PP_ALIGN.CENTER,ls=1.3)
T(s,ML,6.5,CW,0.4,[R("↺  발견한 기준은 다시 DISCOVER로 — 탐색할수록 정교해집니다",12,ACCENT,MD)],
  al=PP_ALIGN.CENTER,an=MSO_ANCHOR.MIDDLE)
foot(s,6)

# ============================================================
# 07~09 · FEATURES
# ============================================================
def feature(idx,num,sec,title,pts,pl,ps):
    s=slide(PAPER); eyebrow(s,num,sec)
    T(s,ML,1.5,6.6,1.7,title,ls=1.2)
    py=3.42; hline(s,ML,py,6.4,c=INK,wt=1.3)
    for h,d in pts:
        T(s,ML,py+0.2,6.4,0.4,[R(h,15,INK,SB)])
        T(s,ML,py+0.62,6.4,0.5,[R(d,12,SUB,RG)],ls=1.3)
        py+=1.05; hline(s,ML,py,6.4,c=HAIR,wt=1.0)
    phone(s,10.35,1.35,4.9,pl,ps); foot(s,idx)

feature(7,"06","RECORD",
  [[R("매일 쓰라고 하지 않습니다.",31,INK,LT,-0.5)],
   [R("7일에 한 번, ",31,INK,LT,-0.5),R("내 말로.",31,ACCENT,SB,-0.5)]],
  [("실제 사용 조합을 그대로 보존","제품·아침/저녁·순서·빈도를 남겨 경험이 생긴 조건을 잃지 않습니다."),
   ("고르고, 내 말로 남기고","‘마음에 들어요·아쉬워요·아직 모르겠어요’ + 자유 원문."),
   ("‘아직 모르겠어요’도 유효한 기록","결론을 강요하지 않아 부담 없이 일상에서 이어집니다.")],
  "홈 · 7일 연구 카드","느낌 남기기 · DAY n/7")

feature(8,"07","EVIDENCE",
  [[R("AI 문장의 새로움이 아니라,",30,INK,LT,-0.5)],
   [R("내 기록이 ",30,INK,LT,-0.5),R("근거",30,ACCENT,SB,-0.5),R("가 됩니다",30,INK,LT,-0.5)]],
  [("지지 근거와 반대 기록을 함께","‘같은 방향 3건 · 다른 기록 1건’처럼 원문을 인용해 보여줍니다."),
   ("환각이 아닌, 검증된 데이터만","서버가 소유권·제품 사실·출력 스키마를 검증한 내 기록만."),
   ("원문은 동의 없이 밖으로 안 나감","자유 메모는 요청하지 않으면 AI 공급자에게 보내지 않습니다.")],
  "AI 채팅 · 근거 인용","‘이 답변에 쓴 근거’ 패널")

feature(9,"08","YOUR MAP",
  [[R("쓸수록 ",31,INK,LT,-0.5),R("나만의",31,ACCENT,SB,-0.5)],
   [R("경험 지도",31,ACCENT,SB,-0.5),R("가 자라납니다",31,INK,LT,-0.5)]],
  [("반복된 경험만 패턴으로","‘최근 기록을 비교해, 반복된 경험만 연결’ — 흐름만 보여줍니다."),
   ("여섯 축은 내 기록으로 동적 생성","길이는 좋고 나쁨이 아니라 연결된 ‘근거량’을 나타냅니다."),
   ("쌓일수록 대체 어려운 개인 자산","발견한 기준은 다음 제품 탐색과 AI 답변에 다시 쓰입니다.")],
  "홈 · 육각 경험 지도","PROFILE · INSIGHT")

# ============================================================
# 10 · DIFFERENTIATION (dark)
# ============================================================
s=slide(INKBG)
eyebrow(s,"09","DIFFERENTIATION",dark=True)
T(s,ML,1.85,CW,2.6,
  [[R("“차별성은 AI가 만든 문장이 아니라,",31,PAPERON,LT,-0.5)],
   [R("한 기록이 다른 제품을 볼 때 ",31,PAPERON,LT,-0.5),R("다시 등장하는 경험",31,SAGEON,SB,-0.5)],
   [R("에서 검증됩니다.”",31,PAPERON,LT,-0.5)]],ls=1.25)
pil=[("내 원본이 자산","남의 후기가 아니라, 내가 남긴 조건과 느낌이 데이터가 된다"),
     ("진단하지 않는다","피부 타입·원인·질환을 판정하지 않아 안전하고 신뢰된다"),
     ("개인 데이터 해자","쌓일수록 다른 서비스가 복제할 수 없는 나만의 지도")]
py=5.15; hline(s,ML,py,CW,c=HAIRON,wt=1.2); cw3=CW/3
for i,(t,d) in enumerate(pil):
    x=ML+i*cw3
    if i: vline(s,x-0.25,py+0.25,1.4,c=HAIRON,wt=1.0)
    T(s,x,py+0.28,cw3-0.5,0.4,[R(f"0{i+1}",14,SAGEON,LT),R("   "+t,15.5,PAPERON,SB)])
    T(s,x,py+0.82,cw3-0.6,0.9,[R(d,11.5,SUBON,RG)],ls=1.35)
foot(s,10,dark=True)

# ============================================================
# 11 · DAILY LIFE
# ============================================================
s=slide(PAPER)
eyebrow(s,"10","DAILY LIFE")
T(s,ML,1.45,11.5,1.3,[R("탐색을 고치지 않고, ",33,INK,LT,-0.5),R("탐색을 자산으로",33,ACCENT,SB,-0.5)])
flow=[("오늘 쓴 느낌","부담 없이 한 줄, 내 말로"),("반복 패턴","AI가 비슷한 경험을 연결"),
      ("다음 탐색에 재등장","제품을 고를 때 내 기준이 먼저")]
fy=3.5; colw=CW/3; hline(s,ML,fy,CW,c=INK,wt=1.3)
for i,(t,d) in enumerate(flow):
    x=ML+i*colw
    if i: vline(s,x-0.05,fy+0.3,1.5,c=HAIR,wt=1.0)
    T(s,x,fy+0.3,colw-0.5,0.6,[R(f"0{i+1}",22,ACCENT,LT)])
    T(s,x,fy+1.0,colw-0.6,0.4,[R(t,18,INK,SB)])
    T(s,x,fy+1.5,colw-0.7,0.4,[R(d,12,SUB,RG)],ls=1.25)
    if i<2: T(s,x+colw-0.55,fy+0.33,0.5,0.6,[R("→",20,FAINT,LT)],al=PP_ALIGN.CENTER)
T(s,ML,5.85,CW,0.9,
  [R("‘써보고 싶다’는 습관을 고치려 하지 않습니다. ",15,INK,RG),
   R("그 습관에서 나온 경험이 사라지지 않게 붙잡아, 다음 선택을 더 쉽게 만드는 것 — 그것이 SKN이 일상에 들어가는 방식입니다.",15,SUB,RG)],ls=1.5)
foot(s,11)

# ============================================================
# 12 · TRACTION
# ============================================================
s=slide(PAPER)
eyebrow(s,"11","TRACTION")
T(s,ML,1.45,7.2,1.0,[R("지금, ",34,INK,LT,-0.5),R("실제로 작동합니다",34,ACCENT,SB,-0.5)])
pts=[("풀스택으로 배포 완료","React 프론트 · 서버 · SQLite. 브라우저에서 바로 체험."),
     ("개인 기록 무결성·소유권 보장","사용자·제품·루틴·경험·패턴의 불변식을 서버가 검증."),
     ("AI 근거 참조를 서버가 검증","AI 출력이 실제 개인 기록·제품 근거만 참조하도록 스키마 검증."),
     ("P0 경험 순환이 한 번 닫힘","탐색→기록→패턴→재사용→Rescue까지 실제 흐름으로 연결.")]
ly=3.05; lw=6.5; hline(s,ML,ly,lw,c=INK,wt=1.3)
for i,(h,d) in enumerate(pts):
    T(s,ML,ly+0.2,0.6,0.72,[R(f"0{i+1}",15,ACCENT,LT)])
    T(s,ML+0.7,ly+0.2,lw-0.7,0.4,[R(h,14.5,INK,SB)])
    T(s,ML+0.7,ly+0.6,lw-0.7,0.5,[R(d,11.5,SUB,RG)],ls=1.25)
    ly+=0.95; hline(s,ML,ly,lw,c=HAIR,wt=1.0)
rx=ML+lw+0.9; vline(s,rx-0.45,3.05,3.4,c=HAIR,wt=1.0)
T(s,rx,3.0,PW-MR-rx,0.4,[R("LIVE",12,ACCENT,SB,3.0)])
phone(s,(rx+PW-MR)/2,3.55,2.5,"데모 QR / 스크린샷","")
T(s,rx,6.2,PW-MR-rx,0.4,[R("skn.today",15,INK,SB,0.3)])
T(s,rx,6.62,PW-MR-rx,0.4,[R("테스트 계정 20개로 즉시 체험",11,SUB,RG)])
foot(s,12)

# ============================================================
# 13 · MARKET
# ============================================================
s=slide(PAPER)
eyebrow(s,"12","MARKET")
T(s,ML,1.45,11.5,1.0,[R("탐색형 소비자가 ",33,INK,LT,-0.5),R("가장 빠르게 크는 시장",33,ACCENT,SB,-0.5)])
pic_w(s,f"{CH}/market.png",ML+0.2,2.75,4.5)
rx=6.2; vline(s,rx-0.5,2.9,3.6,c=HAIR,wt=1.0)
facts=[("11.2조 원","국내 스킨케어 시장 (화장품 24조의 46.8%)"),
       ("연 5.4% 성장","2026–2035 CAGR, 카테고리 내 최고 성장군"),
       ("2030 여성","뷰티 얼리어답터가 초기 진입점 — 탐색 빈도가 가장 높음")]
yy=3.0
for i,(h,d) in enumerate(facts):
    T(s,rx,yy,PW-MR-rx,0.5,[R(h,21,INK,SB,-0.3)])
    T(s,rx,yy+0.5,PW-MR-rx,0.6,[R(d,12,SUB,RG)],ls=1.3)
    yy+=1.2;
    if i<2: hline(s,rx,yy-0.15,PW-MR-rx,c=HAIR,wt=1.0)
source(s,"Expert Market Research 한국 화장품 시장 2025 · 스킨케어 46.8% 점유 · SKN 추정")
foot(s,13)

# ============================================================
# 14 · BUSINESS MODEL (방탄)
# ============================================================
s=slide(PAPER)
eyebrow(s,"13","BUSINESS MODEL")
T(s,ML,1.45,11.5,1.0,[R("쌓일수록 강해지는 ",33,INK,LT,-0.5),R("수익 구조",33,ACCENT,SB,-0.5)])
pic_w(s,f"{CH}/bm.png",ML,2.5,7.0)
# 우: 방어 논리
rx=8.4; vline(s,rx-0.4,2.6,4.1,c=HAIR,wt=1.0)
T(s,rx,2.55,PW-MR-rx,0.4,[R("왜 방어되는가",11,FAINT,SB,1.5)])
defend=[("데이터 해자","개인 경험이 쌓일수록 이탈 비용↑ · 대체 불가"),
        ("검증된 리텐션","재사용·패턴 열람이 반복 방문을 만든다")]
yy=3.0
for h,d in defend:
    T(s,rx,yy,PW-MR-rx,0.4,[R(h,14,INK,SB)])
    T(s,rx,yy+0.38,PW-MR-rx,0.5,[R(d,11,SUB,RG)],ls=1.25)
    yy+=1.0
# 우하: 우리가 하지 않는 것 (선제 방어)
by=5.15
box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(rx),Inches(by),Inches(PW-MR-rx),Inches(1.55))
box.adjustments[0]=0.08; box.fill.solid(); box.fill.fore_color.rgb=PANEL
box.line.fill.background(); box.shadow.inherit=False
T(s,rx+0.3,by+0.22,PW-MR-rx-0.6,0.4,[R("우리가 하지 않는 것",12,ACCENT,SB,0.5)])
T(s,rx+0.3,by+0.62,PW-MR-rx-0.6,0.9,
  [[R("· 개인 기록을 브랜드에 판매",11,INK,MD)],
   [R("· 결과 순위를 돈으로 매수",11,INK,MD)],
   [R("· 개인화 의료광고 · 진단",11,INK,MD)]],ls=1.42)
# 좌하 캡션
T(s,ML,6.75,7.0,0.4,[R("무료로 순환을 검증한 뒤, 구독 → 투명한 제휴로 단계적 확장",11.5,SUB,RG)])
foot(s,14)

# ============================================================
# 15 · CLOSING (dark)
# ============================================================
s=slide(INKBG)
pic_w(s,f"{AS}/orb.png",8.7,1.2,4.3)
pic_w(s,f"{AS}/logo_paper.png",ML,0.82,1.4)
T(s,ML,2.75,11,2.2,
  [[R("결국 남아야 할 것은,",46,PAPERON,LT,-1.0)],
   [R("사용자의 피부와 경험.",46,SAGEON,SB,-1.0)]],ls=1.2)
T(s,ML,5.15,10.5,0.6,[R("써본 만큼 나를 더 잘 알게 되는 스킨케어 경험 아카이브.",14.5,SUBON,RG)])
hline(s,ML,6.5,CW,c=HAIRON,wt=1.0)
T(s,ML,6.68,8,0.35,[R("skn.today",12,SAGEON,MD,0.3)],an=MSO_ANCHOR.MIDDLE)
T(s,PW-MR-4,6.68,4,0.35,[R("AI · WELLNESS",11,SUBON,MD,2.4)],al=PP_ALIGN.RIGHT,an=MSO_ANCHOR.MIDDLE)

out=f"{BASE}/SKN_pitch3.pptx"; prs.save(out)
print("saved:",out,"slides:",len(prs.slides._sldIdLst))
